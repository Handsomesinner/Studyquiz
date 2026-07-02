"""Document processing: extract text from lecture files and split it into chunks.

This is stage (i) of the RAG pipeline described in Chapter 3 of the project
report: text extraction and chunking.

Text-based documents of any kind are accepted. PDF, Word (.docx) and
PowerPoint (.pptx) are parsed with dedicated libraries; everything else is
read as plain text (covers .txt, .md, .csv, .html, source files, …).
Scanned/image-only files still can't be read — OCR is outside the project
scope — so genuinely binary uploads are rejected with a clear message.
"""

import io
import re

from pypdf import PdfReader

# Chunks are sized in words. Roughly 200 words per chunk keeps each chunk
# focused on one idea, and the overlap stops a sentence that straddles a
# boundary from being lost to both chunks.
CHUNK_SIZE_WORDS = 200
CHUNK_OVERLAP_WORDS = 40


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs]
    # Table cells hold text too — lecture handouts often use tables.
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_plain_text(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    # Guard against binary files (images, video, archives) reaching here via
    # the "accept anything" path: if the decode is mostly replacement or
    # control characters, there is no real text to work with.
    if text:
        bad = sum(
            1 for c in text if c == "�" or (ord(c) < 32 and c not in "\t\n\r")
        )
        if bad / len(text) > 0.30:
            raise ValueError(
                "This file doesn't appear to contain readable text. Scanned "
                "images and media files aren't supported (OCR is out of "
                "scope). Try a PDF, Word, PowerPoint, or text document."
            )
    return text


def extract_text(filename: str, data: bytes) -> str:
    """Extract raw text from an uploaded document of any text-based type."""
    name = filename.lower()
    if name.endswith(".pdf"):
        text = _extract_pdf(data)
    elif name.endswith(".docx"):
        text = _extract_docx(data)
    elif name.endswith(".pptx"):
        text = _extract_pptx(data)
    else:
        text = _extract_plain_text(data)

    # Collapse repeated whitespace that document extraction tends to produce.
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str) -> list[str]:
    """Split text into overlapping word-window chunks.

    A sliding window is the simplest chunking strategy that still preserves
    local context; more advanced strategies (per-heading, per-paragraph)
    are discussed as future work in the report.
    """
    words = text.split()
    if not words:
        return []

    chunks = []
    step = CHUNK_SIZE_WORDS - CHUNK_OVERLAP_WORDS
    for start in range(0, len(words), step):
        window = words[start : start + CHUNK_SIZE_WORDS]
        if len(window) < 30 and chunks:
            # Tail too small to be a useful chunk on its own — merge into
            # the previous one instead of creating a fragment.
            chunks[-1] = chunks[-1] + " " + " ".join(window)
            break
        chunks.append(" ".join(window))
    return chunks
